"""PPTX slide text extraction via python-pptx (speaker notes omitted)."""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from app.parsers._versions import package_version
from app.parsers.base import DocumentParser, ParserResult
from app.services.text_extraction_service import normalize_extension


class PptxParser(DocumentParser):
    PARSER_NAME = "pptx_parser"

    def __init__(self) -> None:
        self._version = package_version("python-pptx")

    def supports(self, extension: str, mime_type: str | None) -> bool:
        return normalize_extension(extension) == "pptx"

    def _lines_from_shape(self, shape) -> list[str]:  # type: ignore[no-untyped-def]
        lines: list[str] = []
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            try:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
            except Exception:
                pass
        if getattr(shape, "has_table", False) and shape.has_table:
            try:
                for row in shape.table.rows:
                    cells = [
                        (cell.text or "").replace("\n", " ").strip()
                        for cell in row.cells
                    ]
                    if any(cells):
                        lines.append(" | ".join(cells))
            except Exception:
                pass
        return lines

    def parse_bytes(
        self, content: bytes, filename: str, extension: str
    ) -> ParserResult:
        try:
            prs = Presentation(BytesIO(content))
        except Exception:
            return ParserResult(
                success=False,
                extracted_text=None,
                text_length=0,
                parser_name=self.PARSER_NAME,
                parser_version=self._version,
                metadata={},
                error_code="PARSING_FAILED",
                error_message="Document parsing failed",
            )

        parts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_bits: list[str] = []
            for shape in slide.shapes:
                slide_bits.extend(self._lines_from_shape(shape))
            if slide_bits:
                parts.append(f"--- slide {i} ---\n" + "\n".join(slide_bits))

        text = "\n\n".join(parts).strip()
        return ParserResult(
            success=True,
            extracted_text=text,
            text_length=len(text),
            parser_name=self.PARSER_NAME,
            parser_version=self._version,
            metadata={"slide_count": len(prs.slides)},
        )


__all__ = ["PptxParser"]
